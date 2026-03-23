from __future__ import annotations

import asyncio
import hashlib
import time
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

from docx import Document as DocxReader
from fastapi import Depends, UploadFile
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.celery_core import celery_app
from app.core.milvus_core import vector_manager
from app.core.storage import (get_storage_backend,
                              get_storage_backend_for_provider,
                              resolve_local_parsed_dir)
from app.models import Resources
from app.service.user_service import get_db


class ResourceService:
    _MILVUS_TEXT_MAX_LEN = 2048
    _MILVUS_PARENT_TEXT_MAX_LEN = 4096
    _SPLIT_CHUNK_SIZE = 1800
    _SPLIT_CHUNK_OVERLAP = 200
    _PARENT_SPLIT_CHUNK_SIZE = 3600
    _PARENT_SPLIT_CHUNK_OVERLAP = 300

    def __init__(self, session: AsyncSession):
        self._session = session

    async def create_from_upload(
        self,
        *,
        uid: str,
        file: UploadFile,
        rname: Optional[str] = None,
        rtype: str = "doc",
    ) -> Resources:
        storage = get_storage_backend()
        # File I/O is blocking; offload to a worker thread to keep the event loop responsive.
        stored = await asyncio.to_thread(storage.save_file, file.file, file.filename)

        resource = Resources(
            rid=str(uuid.uuid4()),
            uid=uid,
            rname=rname or (file.filename or stored.key),
            rtype=rtype,
            storage_provider=stored.provider,
            storage_key=stored.key,
            # Keep legacy url populated for now; API derives url at read-time.
            url=stored.url,
        )
        self._session.add(resource)
        await self._session.flush()
        await self._session.refresh(resource)
        return resource

    async def list_by_uid(self, uid: str) -> list[Resources]:
        result = await self._session.execute(
            select(Resources)
            .where(Resources.uid == uid)
            .order_by(Resources.created_at.desc())
        )
        return list(result.scalars().all())

    async def get_by_rid(self, rid: str) -> Optional[Resources]:
        result = await self._session.execute(
            select(Resources).where(Resources.rid == rid)
        )
        return result.scalar_one_or_none()

    @staticmethod
    def resolve_url(resource: Resources) -> str:
        """
        Public URL resolution rule:
        1) If storage_provider + storage_key exists: derive via backend.get_url(key)
        2) Else fall back to legacy url (maybe external)
        """
        if resource.storage_provider and resource.storage_key:
            backend = get_storage_backend_for_provider(resource.storage_provider)
            return backend.get_url(resource.storage_key)
        return resource.url or ""

    async def parse_resource(self, *, rid: str, force: bool = False) -> dict[str, Any]:
        """
        Enqueue a resource parsing task to Celery.
        """
        result = await self._session.execute(
            select(Resources).where(Resources.rid == rid).with_for_update()
        )
        resource = result.scalar_one_or_none()
        if resource is None:
            raise ValueError("Resource not found")

        # Prevent repeated enqueue when users click execute multiple times.
        if resource.status == 1:
            return {
                "rid": rid,
                "task_id": None,
                "status": "parsing",
                "message": "资源正在解析中，请勿重复提交。",
            }
        if resource.status == 2 and not force:
            return {
                "rid": rid,
                "task_id": None,
                "status": "parsed",
                "message": "资源已解析完成。如需重跑请使用 force=true。",
            }

        # Mark as parsing before enqueue to minimize duplicate tasks under concurrency.
        resource.status = 1
        await self._session.flush()

        payload = {"rid": rid, "uid": resource.uid}
        async_result = celery_app.send_task("chatmooc.parse_resource", args=[payload])
        return {
            "rid": rid,
            "task_id": async_result.id,
            "status": "queued",
            "message": "解析任务已创建。",
        }

    async def parse_resource_now(self, *, rid: str) -> dict[str, Any]:
        """
        1) 提取文本信息并保存为 markdown 文件
        2) 预留摘要和关键词接口（当前占位实现）
        3) 将文本分块后写入 Milvus
        """
        resource = await self.get_by_rid(rid)
        if resource is None:
            raise ValueError("Resource not found")

        resource.status = 1  # parsing
        await self._session.flush()

        raw_bytes = await asyncio.to_thread(self._load_resource_bytes, resource)
        extracted_text = await asyncio.to_thread(
            self._extract_text_from_bytes,
            raw_bytes,
            resource.rname,
            resource.rtype,
        )
        normalized_text = (extracted_text or "").strip()
        if not normalized_text:
            resource.status = 0
            await self._session.flush()
            raise ValueError("No extractable text content found in resource")

        await asyncio.to_thread(self._write_parsed_markdown, resource, normalized_text)
        resource.content = None
        summary, keywords = await self._generate_summary_and_keywords(normalized_text)
        resource.summary = summary
        resource.keywords = keywords

        chunk_entries = self._split_text_with_parent(normalized_text)
        if not chunk_entries:
            resource.status = 0
            await self._session.flush()
            raise ValueError("Text splitting produced no chunks")

        created_ts = int(time.time())
        docs = [
            Document(
                page_content=entry["text"],
                metadata={
                    "resource_id": resource.rid,
                    "user_id": resource.uid,
                    "parent_text": entry["parent_text"],
                    "create_time": created_ts,
                },
            )
            for entry in chunk_entries
        ]
        ids = [
            self._build_chunk_pk(
                rid=resource.rid,
                index=index,
                content=doc.page_content,
            )
            for index, doc in enumerate(docs)
        ]

        vector_store = vector_manager.store
        if vector_store is None:
            resource.status = 0
            await self._session.flush()
            raise RuntimeError("Knowledge base (Milvus) is not initialized")

        # Idempotency: replace existing vectors for the same resource before inserting.
        filter_expr = f'resource_id == "{resource.rid}" and user_id == "{resource.uid}"'
        if hasattr(vector_store, "adelete"):
            await vector_store.adelete(expr=filter_expr)
        else:
            await asyncio.to_thread(vector_store.delete, expr=filter_expr)

        # Collection pk is required (varchar); pass explicit ids for primary_field="pk".
        await vector_store.aadd_documents(docs, ids=ids)

        resource.status = 2  # parsed
        await self._session.flush()
        await self._session.refresh(resource)
        return {
            "rid": resource.rid,
            "status": resource.status,
            "chunks": len(chunk_entries),
        }

    def _load_resource_bytes(self, resource: Resources) -> bytes:
        provider = (resource.storage_provider or "").strip().lower()
        if provider == "local":
            backend = get_storage_backend_for_provider(provider)
            base_dir = getattr(backend, "base_dir", None)
            if base_dir is None:
                raise RuntimeError("Local storage backend is missing base_dir")
            if not resource.storage_key:
                raise ValueError("Resource storage_key is empty")
            file_path = Path(base_dir) / resource.storage_key
            if not file_path.exists():
                raise FileNotFoundError(
                    f"Resource file not found: {resource.storage_key}"
                )
            return file_path.read_bytes()
        raise NotImplementedError(
            f"Storage provider {provider or '<empty>'} not supported for parsing"
        )

    def _write_parsed_markdown(self, resource: Resources, text: str) -> Path:
        parsed_dir = resolve_local_parsed_dir()
        filename = f"{resource.rid}.md"
        file_path = parsed_dir / filename
        file_path.write_text(text, encoding="utf-8")
        return file_path

    def _extract_text_from_bytes(
        self, payload: bytes, filename: str, rtype: Optional[str]
    ) -> str:
        type_hint = (rtype or "").strip().lower()
        suffix = Path(filename or "").suffix.lower()

        if type_hint == "video":
            raise NotImplementedError("Video resource parsing is not implemented yet")
        if type_hint == "audio":
            raise NotImplementedError("Audio resource parsing is not implemented yet")

        if suffix == ".pdf":
            return self._extract_pdf_text(payload)
        if suffix == ".docx":
            return self._extract_docx_text(payload)
        if suffix == ".pptx":
            return self._extract_pptx_text(payload)
        if suffix in {".txt", ".csv", ".log"}:
            return payload.decode("utf-8", errors="ignore")
        if suffix in {".md", ".markdown"}:
            return payload.decode("utf-8", errors="ignore")

        if type_hint == "doc":
            return payload.decode("utf-8", errors="ignore")

        # Unknown extensions fallback to plain text decode.
        return payload.decode("utf-8", errors="ignore")

    @staticmethod
    def _extract_pdf_text(payload: bytes) -> str:
        reader = PdfReader(BytesIO(payload))
        pages: list[str] = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        return "\n".join(pages)

    @staticmethod
    def _extract_docx_text(payload: bytes) -> str:
        document = DocxReader(BytesIO(payload))
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        return "\n".join(lines)

    @staticmethod
    def _extract_pptx_text(payload: bytes) -> str:
        presentation = Presentation(BytesIO(payload))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    lines.append(text)
        return "\n".join(lines)

    @staticmethod
    def _split_text(text: str) -> list[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=ResourceService._SPLIT_CHUNK_SIZE,
            chunk_overlap=ResourceService._SPLIT_CHUNK_OVERLAP,
            separators=["\n\n", "\n", "。", "，", " ", ""],
        )
        chunks = splitter.split_text(text)
        # Guarantee every chunk fits Milvus `text` varchar(2048) in UTF-8 bytes.
        normalized: list[str] = []
        for chunk in chunks:
            if len(chunk.encode("utf-8")) <= ResourceService._MILVUS_TEXT_MAX_LEN:
                normalized.append(chunk)
                continue
            normalized.extend(
                ResourceService._split_by_utf8_bytes(
                    chunk,
                    max_bytes=ResourceService._MILVUS_TEXT_MAX_LEN,
                    overlap_chars=ResourceService._SPLIT_CHUNK_OVERLAP,
                )
            )
        return normalized

    @staticmethod
    def _split_text_with_parent(text: str) -> list[dict[str, str]]:
        """
        Hierarchical split:
        - parent_text: larger context window (<= 4096 bytes)
        - text: child chunk for embedding/indexing (<= 2048 bytes)
        """
        parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=ResourceService._PARENT_SPLIT_CHUNK_SIZE,
            chunk_overlap=ResourceService._PARENT_SPLIT_CHUNK_OVERLAP,
            separators=["\n\n", "\n", "。", "，", " ", ""],
        )
        parent_chunks = parent_splitter.split_text(text)

        entries: list[dict[str, str]] = []
        for parent in parent_chunks:
            parent_safe = ResourceService._truncate_utf8_bytes(
                parent,
                ResourceService._MILVUS_PARENT_TEXT_MAX_LEN,
            )
            child_chunks = ResourceService._split_text(parent_safe)
            for child in child_chunks:
                entries.append(
                    {
                        "text": child,
                        "parent_text": parent_safe,
                    }
                )
        return entries

    @staticmethod
    def _truncate_utf8_bytes(text: str, max_bytes: int) -> str:
        encoded = text.encode("utf-8")
        if len(encoded) <= max_bytes:
            return text
        return encoded[:max_bytes].decode("utf-8", errors="ignore")

    @staticmethod
    def _split_by_utf8_bytes(
        text: str, max_bytes: int, overlap_chars: int
    ) -> list[str]:
        pieces: list[str] = []
        cursor = 0
        text_len = len(text)
        overlap = max(0, overlap_chars)

        while cursor < text_len:
            low = cursor + 1
            high = text_len
            best_end = cursor + 1

            # Binary search max end index that still fits utf-8 bytes length.
            while low <= high:
                mid = (low + high) // 2
                candidate = text[cursor:mid]
                if len(candidate.encode("utf-8")) <= max_bytes:
                    best_end = mid
                    low = mid + 1
                else:
                    high = mid - 1

            piece = text[cursor:best_end]
            if not piece:
                # Fallback safety for any pathological case.
                piece = ResourceService._truncate_utf8_bytes(
                    text[cursor : cursor + 1], max_bytes
                )
                best_end = cursor + max(1, len(piece))

            pieces.append(piece)
            if best_end >= text_len:
                break

            cursor = max(best_end - overlap, cursor + 1)

        return pieces

    @staticmethod
    def _build_chunk_pk(*, rid: str, index: int, content: str) -> str:
        digest = hashlib.sha1(content.encode("utf-8")).hexdigest()[:12]
        # Keep length <= 64 to match Milvus pk varchar(64).
        base = f"{rid.replace('-', '')[:24]}-{index:04d}-{digest}"
        return base[:64]

    @staticmethod
    async def _generate_summary_and_keywords(
        text: str,
    ) -> tuple[Optional[str], Optional[dict[str, Any]]]:
        """
        Placeholder for future LLM pipeline.
        """
        _ = text
        return None, None


def get_resource_service(db: AsyncSession = Depends(get_db)) -> ResourceService:
    return ResourceService(db)
